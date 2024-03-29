import csv 
import os
import string
import xml.etree.ElementTree as ET 
from pptx import Presentation
from pptx.util import Inches
import html2text
import datetime
import re
from tkinter import filedialog
from tkinter import *

#remove empty lines from xml-sourced text
def stripEmptyLines(text):
    return os.linesep.join([s for s in text.splitlines() if s])

#check if this is a folder cpntaining module.xml
def checkforxml(relpath):
    for file in os.listdir(relpath):
        if file == 'module.xml':
            return True;
    return False

#parse XMl file into scenes/slides and identify presentation title
def parseXML(xmlfile): 
  
    # create element tree object 
    tree = ET.parse(xmlfile) 
    root = tree.getroot() 
    sceneItems = []

    for item in root.findall('./lesson/scenes/scene'): 
        sceneItems.append(item)     
    return sceneItems

def createPresentation(folder):
    templatePath = 'MPA2019Template.pptx'
    prs = Presentation(templatePath)
    slide_layout = prs.slide_layouts[1]
    xmlpath = folder + '/module.xml'
    xmlScenes = parseXML(xmlpath)
    slideCount = 0

    for scene in xmlScenes:
        if slideCount == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]
        slideCount += 1

        titleText = scene.find('sceneTitle').text
        contentText = html2text.html2text(scene.find('content').text)
        contentText = contentText.replace('* ', '')
        contentText = stripEmptyLines(contentText)

        #create slide
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = titleText
        tf = body_shape.text_frame
        tf.text = contentText
       
        #add images if anys
        for item in scene.findall('.sceneAssets/asset'):
            if item.text is not None:
                if (contentText.count('\n') == 0 and contentText.endswith('.png')):
                    left = Inches(6)
                    height = Inches(6)
                    top = Inches(0.7)
                    tf.text = ''
                else:
                    left = Inches(0.5)
                    height = Inches(2)
                    top = Inches(0.4)
                img_path = os.path.join(folder, item.text)
                pic = shapes.add_picture(img_path, left, top, height = height)

        #add cue point times as reference notes inside the powerpoint file               
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        cueTime = str(datetime.timedelta(seconds = float(scene.attrib['cuePoint'])))
        print(cueTime)
        text_frame.text = ("cue time: " + cueTime)

    prs.save(folder +'.pptx')
    print(folder + "exported " + str(slideCount) + " slides")
     
def main():
	root = Tk()
	root.withdraw()
	folder_selected = filedialog.askdirectory() 
	for entry in os.scandir(folder_selected):
		if entry.is_dir():
			relpath = os.path.relpath(entry.path)
			print(relpath)
			if checkforxml(relpath):
				print("checked!: " + relpath)
				createPresentation(relpath)
			else:
				print("no xml module.xml found in " + relpath)
      
if __name__ == "__main__": 
  
    # calling main function 
    main() 